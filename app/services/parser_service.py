"""
services/parser_service.py — 文件解析層

架構說明（借鑑 RAG-Anything 設計理念）
═══════════════════════════════════════════════════════
核心洞見：所有格式統一歸一到 content_list，再由同一套邏輯處理。

解析引擎（由 config.PDF_PARSER 控制）
──────────────────────────────────────
"auto"       → docling → pdfplumber（依序嘗試）
"docling"    → 只用 Docling
"paddleocr"  → 已停用（PaddleX 3.x 需連網），fallback 到 pdfplumber
"pdfplumber" → 只用 pdfplumber

各格式處理策略
──────────────
PDF    → 引擎直接解析
Office → Docling 原生 → LibreOffice/win32com 轉 PDF → 格式專用 fallback
TXT/MD → ReportLab 轉 PDF → 引擎解析 / 直接讀文字
Image  → Ollama Vision

content_list 統一格式（每個元素）
──────────────────────────────────
  type          : "text" | "image" | "table" | "equation"
  page_idx      : 0-based 頁碼
  text          : 文字內容
  text_level    : 標題層級（0=正文，1+=標題）
  img_path      : 圖片絕對路徑
  image_caption : 圖片說明
  table_body    : 表格內容
  table_caption : 表格標題

Vision context window
─────────────────────
描述圖片時帶入前後 VISION_CONTEXT_PAGES 頁文字作為 context。

Bug fixes（v2）
───────────────
- Docling：設定 DOCLING_OFFLINE=1 + HF_DATASETS_OFFLINE=1 避免
  啟動時嘗試下載 modelscope / HuggingFace 模型導致失敗
- PaddleOCR：v3+ 移除 show_log 參數，改用 logging level 控制；
  初始化時多重嘗試，逐步移除不相容參數

不含 chunking、embedding、DB 邏輯。
"""
from __future__ import annotations

import base64
import hashlib
import io
import json
import logging
import os
import re
import shutil
import subprocess
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

import pdfplumber

from app.core.config import settings
from app.models.schemas import ContentType
from app.services.model_service import describe_image
from app.utils.helpers import get_logger

logger = get_logger(__name__)


# ══════════════════════════════════════════════════════════
# 資料結構
# ══════════════════════════════════════════════════════════

@dataclass
class ParsedPage:
    """一個「邏輯頁」的解析結果"""
    page_num     : int
    text         : str
    tables       : list[str]   = field(default_factory=list)
    figure_desc  : list[str]   = field(default_factory=list)
    content_type : ContentType = ContentType.MIXED
    _raw_content_list : list[dict] = field(default_factory=list, repr=False)

    def to_full_text(self) -> str:
        parts = [self.text.strip()]
        for tbl in self.tables:
            parts.append(f"\n\n{tbl}")
        for desc in self.figure_desc:
            parts.append(f"\n\n> **圖片描述**\n{desc}")
        return "\n".join(p for p in parts if p.strip())


# ══════════════════════════════════════════════════════════
# 共用工具
# ══════════════════════════════════════════════════════════

def _md5(data: bytes) -> str:
    return hashlib.md5(data).hexdigest()


def _clean_cell(cell) -> str:
    if cell is None:
        return ""
    text = str(cell).strip()
    text = re.sub(r"\s*\n\s*", " ", text)
    text = text.replace("|", "\\|")
    return text


def _table_to_md(rows: list) -> str:
    if not rows:
        return ""
    cleaned = [[_clean_cell(c) for c in row] for row in rows]
    valid   = [r for r in cleaned if any(r)]
    if len(valid) < 2:
        return ""
    max_c = max(len(r) for r in valid)
    for r in valid:
        r += [""] * (max_c - len(r))
    if all(all(c == "" for c in r) for r in valid):
        return ""
    lines = [
        "| " + " | ".join(valid[0]) + " |",
        "| " + " | ".join(["---"] * max_c) + " |",
    ]
    for row in valid[1:]:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def _table_body_to_md(body: Any) -> str:
    if isinstance(body, str):
        return body.strip()
    if isinstance(body, list):
        if not body:
            return ""
        if all(isinstance(r, (list, tuple)) for r in body):
            return _table_to_md(body)
        return "\n".join(str(r) for r in body)
    return str(body).strip() if body else ""


def _is_scanned_page(text: str) -> bool:
    if len(text) < settings.PDF_SCAN_MIN_CHARS:
        return True
    readable = len(re.findall(r"[\w\u4e00-\u9fff]", text))
    return (readable / len(text)) < settings.PDF_SCAN_MIN_RATIO if text else True


def _safe_read_image_ref(img_obj) -> bytes | None:
    try:
        data = img_obj._data()
        if isinstance(data, (bytes, bytearray)) and len(data) >= settings.IMG_MIN_BYTES:
            return bytes(data)
    except Exception:
        pass
    ref = getattr(img_obj, "ref", None)
    if ref is None:
        return None
    try:
        if hasattr(ref, "read"):
            ref.seek(0)
            data = ref.read()
        elif isinstance(ref, (bytes, bytearray)):
            data = bytes(ref)
        else:
            return None
        return data if len(data) >= settings.IMG_MIN_BYTES else None
    except Exception:
        return None


# ══════════════════════════════════════════════════════════
# Vision context window
# ══════════════════════════════════════════════════════════

def _build_context_from_content_list(
    content_list : list[dict],
    current_page : int,
) -> str:
    window    = settings.VISION_CONTEXT_PAGES
    max_chars = settings.VISION_CONTEXT_MAX_CHARS
    if window == 0 or not content_list:
        return ""
    start = max(0, current_page - window)
    end   = current_page + window + 1
    parts = [
        f"[第 {item.get('page_idx', 0) + 1} 頁] {str(item.get('text', '')).strip()}"
        for item in content_list
        if item.get("type") == "text"
        and item.get("page_idx", 0) != current_page
        and start <= item.get("page_idx", 0) < end
        and str(item.get("text", "")).strip()
    ]
    ctx = "\n".join(parts)
    return (ctx[:max_chars] + "…") if len(ctx) > max_chars else ctx


def _build_context_from_texts(all_texts: list[str], current_idx: int) -> str:
    window    = settings.VISION_CONTEXT_PAGES
    max_chars = settings.VISION_CONTEXT_MAX_CHARS
    if window == 0 or not all_texts:
        return ""
    start = max(0, current_idx - window)
    end   = min(len(all_texts), current_idx + window + 1)
    parts = [
        f"[第 {i + 1} 頁] {all_texts[i].strip()}"
        for i in range(start, end)
        if i != current_idx and all_texts[i].strip()
    ]
    ctx = "\n".join(parts)
    return (ctx[:max_chars] + "…") if len(ctx) > max_chars else ctx


def _describe_with_context(
    img_bytes: bytes, page_num: int, context: str, extra_caption: str = ""
) -> str:
    caption_part = f"\n【圖片說明】{extra_caption}" if extra_caption else ""
    if context:
        full_ctx = (
            f"第 {page_num} 頁圖片{caption_part}\n\n"
            f"【周圍文字 context（前後 {settings.VISION_CONTEXT_PAGES} 頁）】\n"
            f"{context}"
        )
    else:
        full_ctx = f"第 {page_num} 頁圖片{caption_part}"
    return describe_image(img_bytes, context=full_ctx)


# ══════════════════════════════════════════════════════════
# 格式轉換工具
# ══════════════════════════════════════════════════════════

def _find_libreoffice() -> str | None:
    candidates = [
        "libreoffice", "soffice",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "/usr/bin/libreoffice", "/usr/bin/soffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for c in candidates:
        if shutil.which(c):
            return c
        if os.path.isfile(c):
            return c
    return None


def _office_to_pdf_libreoffice(src_bytes: bytes, src_suffix: str) -> bytes | None:
    lo = _find_libreoffice()
    if not lo:
        return None
    with tempfile.TemporaryDirectory() as tmp:
        src = Path(tmp) / f"input{src_suffix}"
        src.write_bytes(src_bytes)
        try:
            r = subprocess.run(
                [lo, "--headless", "--convert-to", "pdf", "--outdir", tmp, str(src)],
                capture_output=True, timeout=60,
            )
            if r.returncode != 0:
                logger.warning("LibreOffice 轉 PDF 失敗：%s",
                               r.stderr.decode(errors="replace")[:200])
                return None
            pdfs = list(Path(tmp).glob("*.pdf"))
            data = pdfs[0].read_bytes() if pdfs else None
            return data if data and len(data) > 100 else None
        except subprocess.TimeoutExpired:
            logger.warning("LibreOffice 逾時")
            return None
        except Exception as e:
            logger.warning("LibreOffice 失敗：%s", e)
            return None


def _office_to_pdf_win32com(src_bytes: bytes, src_suffix: str) -> bytes | None:
    if sys.platform != "win32":
        return None
    try:
        import win32com.client
    except ImportError:
        return None
    fmt_map = {
        ".doc": ("Word.Application", 17), ".docx": ("Word.Application", 17),
        ".xls": ("Excel.Application", 57), ".xlsx": ("Excel.Application", 57),
        ".ppt": ("PowerPoint.Application", 32), ".pptx": ("PowerPoint.Application", 32),
    }
    key = src_suffix.lower()
    if key not in fmt_map:
        return None
    app_name, fmt_code = fmt_map[key]
    with tempfile.TemporaryDirectory() as tmp:
        src_path = os.path.join(tmp, f"input{src_suffix}")
        out_path = os.path.join(tmp, "output.pdf")
        with open(src_path, "wb") as f:
            f.write(src_bytes)
        app = obj = None
        try:
            app = win32com.client.DispatchEx(app_name)
            app.Visible = False
            if hasattr(app, "DisplayAlerts"):
                app.DisplayAlerts = False
            abs_src, abs_out = os.path.abspath(src_path), os.path.abspath(out_path)
            if app_name == "Word.Application":
                obj = app.Documents.Open(abs_src)
                obj.SaveAs(abs_out, FileFormat=fmt_code)
            elif app_name == "Excel.Application":
                obj = app.Workbooks.Open(abs_src, ReadOnly=True)
                obj.ExportAsFixedFormat(0, abs_out)
            elif app_name == "PowerPoint.Application":
                obj = app.Presentations.Open(abs_src, ReadOnly=True, WithWindow=False)
                obj.SaveAs(abs_out, fmt_code)
            return open(out_path, "rb").read() if os.path.exists(out_path) else None
        except Exception as e:
            logger.warning("win32com 轉 PDF 失敗：%s", e)
            return None
        finally:
            import time
            for o in (obj, app):
                if o is None:
                    continue
                for method in ("Close", "Quit"):
                    try:
                        getattr(o, method)(False) if method == "Close" else getattr(o, method)()
                        break
                    except Exception:
                        pass
            time.sleep(0.3)


def _office_to_pdf(src_bytes: bytes, src_suffix: str) -> tuple[bytes, bool]:
    for converter in (_office_to_pdf_libreoffice, _office_to_pdf_win32com):
        result = converter(src_bytes, src_suffix)
        if result:
            logger.info("%s 轉 PDF 成功（%s）", src_suffix, converter.__name__)
            return result, True
    logger.warning("無法將 %s 轉為 PDF", src_suffix)
    return src_bytes, False


def _txt_to_pdf(file_bytes: bytes) -> bytes | None:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import inch
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    except ImportError:
        return None
    try:
        text = file_bytes.decode("utf-8", errors="replace")
        buf  = io.BytesIO()
        doc  = SimpleDocTemplate(buf, pagesize=A4,
                                 leftMargin=inch, rightMargin=inch,
                                 topMargin=inch, bottomMargin=inch)
        try:
            pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
            font_name = "STSong-Light"
        except Exception:
            font_name = "Helvetica"
        styles        = getSampleStyleSheet()
        normal        = styles["Normal"]
        normal.fontName = font_name
        story = []
        for line in text.splitlines():
            s = line.strip()
            if not s:
                story.append(Spacer(1, 6))
                continue
            safe = s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
            story.append(Paragraph(safe, normal))
            story.append(Spacer(1, 3))
        if not story:
            story.append(Paragraph("(empty)", normal))
        doc.build(story)
        return buf.getvalue()
    except Exception as e:
        logger.warning("TXT → PDF 失敗：%s", e)
        return None


# ══════════════════════════════════════════════════════════
# content_list → ParsedPage list（統一出口）
# ══════════════════════════════════════════════════════════

def _content_list_to_pages(
    content_list : list[dict],
    images_base  : Path | None = None,
    seen_img     : dict[str, str] | None = None,
) -> list[ParsedPage]:
    """所有引擎的統一出口，圖片 context window、表格、公式只寫一次"""
    if seen_img is None:
        seen_img = {}

    from collections import defaultdict
    by_page: dict[int, list[dict]] = defaultdict(list)
    for item in content_list:
        by_page[item.get("page_idx", 0)].append(item)

    if not by_page:
        return []

    total_pages = max(by_page.keys()) + 1
    pages = [
        ParsedPage(page_num=i + 1, text="", _raw_content_list=content_list)
        for i in range(total_pages)
    ]

    for page_idx in range(total_pages):
        items      = by_page.get(page_idx, [])
        text_parts : list[str] = []
        tbl_mds    : list[str] = []
        fig_descs  : list[str] = []
        eq_parts   : list[str] = []

        for item in items:
            typ = item.get("type", "")

            if typ == "text":
                t     = str(item.get("text", "")).strip()
                level = item.get("text_level", 0) or 0
                if t:
                    text_parts.append("#" * level + " " + t if level > 0 else t)

            elif typ == "equation":
                eq = str(item.get("text", "") or item.get("latex", "")).strip()
                if eq:
                    eq_parts.append(f"$$\n{eq}\n$$")

            elif typ == "table":
                body     = item.get("table_body") or item.get("table_data") or item.get("text", "")
                captions = item.get("table_caption", "")
                caption  = (" ".join(captions) if isinstance(captions, list)
                            else str(captions)).strip()
                md = _table_body_to_md(body)
                if md:
                    tbl_mds.append((f"**{caption}**\n\n" if caption else "") + md)

            elif typ == "image":
                img_path_str = item.get("img_path", "")
                captions     = item.get("image_caption", item.get("img_caption", []))
                caption_text = (" ".join(captions) if isinstance(captions, list)
                               else str(captions)).strip()

                img_bytes_data: bytes | None = None
                if img_path_str:
                    img_path = Path(img_path_str)
                    if not img_path.is_absolute() and images_base:
                        img_path = images_base / img_path_str
                    if img_path.exists():
                        try:
                            img_bytes_data = img_path.read_bytes()
                        except Exception:
                            pass

                if img_bytes_data is None:
                    uri = item.get("image_uri", item.get("img_uri", ""))
                    if uri and uri.startswith("data:"):
                        try:
                            _, b64 = uri.split(",", 1)
                            img_bytes_data = base64.b64decode(b64)
                        except Exception:
                            pass

                if img_bytes_data and len(img_bytes_data) >= settings.IMG_MIN_BYTES:
                    h = _md5(img_bytes_data)
                    if h not in seen_img:
                        ctx  = _build_context_from_content_list(content_list, page_idx)
                        desc = _describe_with_context(
                            img_bytes_data, page_idx + 1, ctx, caption_text
                        )
                        seen_img[h] = desc
                    fig_descs.append(seen_img[h])

        full_text = "\n".join(text_parts)
        if eq_parts:
            full_text = full_text + "\n\n" + "\n\n".join(eq_parts)
        full_text = re.sub(r"\n{3,}", "\n\n", full_text).strip()

        if fig_descs and not full_text and not tbl_mds:
            ctype = ContentType.IMAGE
        elif tbl_mds and not full_text and not fig_descs:
            ctype = ContentType.TABLE
        elif eq_parts and not tbl_mds and not fig_descs:
            ctype = ContentType.EQUATION
        else:
            ctype = ContentType.MIXED

        pages[page_idx].text         = full_text
        pages[page_idx].tables       = tbl_mds
        pages[page_idx].figure_desc  = fig_descs
        pages[page_idx].content_type = ctype

    return [p for p in pages if p.to_full_text().strip()]


# ══════════════════════════════════════════════════════════
# Docling 引擎
# ══════════════════════════════════════════════════════════

def _check_docling() -> bool:
    try:
        from docling.document_converter import DocumentConverter  # noqa
        return True
    except ImportError:
        return False


def _run_docling(file_bytes: bytes, suffix: str) -> list[dict] | None:
    """
    Docling Python API 解析。
    Bug fix：設定離線環境變數，避免啟動時嘗試下載 modelscope/HuggingFace 模型。
    """
    if not _check_docling():
        return None
    try:
        from docling.document_converter import DocumentConverter, PdfFormatOption
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import PdfPipelineOptions
    except ImportError:
        return None

    with tempfile.TemporaryDirectory() as tmp:
        src     = Path(tmp) / f"input{suffix}"
        img_dir = Path(tmp) / "images"
        src.write_bytes(file_bytes)
        img_dir.mkdir()

        # 在 import 之前設定離線環境變數（必須在 import 前設定才有效）
        # 這裡補設是為了確保子執行緒也繼承到
        os.environ["DOCLING_OFFLINE"]       = "1"
        os.environ["HF_DATASETS_OFFLINE"]   = "1"
        os.environ["TRANSFORMERS_OFFLINE"]  = "1"
        os.environ["HF_HUB_OFFLINE"]        = "1"   # huggingface_hub 離線
        os.environ["CURL_CA_BUNDLE"]        = ""    # 停用 SSL 憑證驗證（自簽憑證環境）
        os.environ["REQUESTS_CA_BUNDLE"]    = ""

        try:
            pipeline_options = PdfPipelineOptions()
            if hasattr(pipeline_options, "generate_picture_images"):
                pipeline_options.generate_picture_images = True
            if hasattr(pipeline_options, "images_scale"):
                pipeline_options.images_scale = 2.0
            # 關閉 OCR 和表格結構模型，避免觸發連網下載
            if hasattr(pipeline_options, "do_ocr"):
                pipeline_options.do_ocr = False
            if hasattr(pipeline_options, "do_table_structure"):
                pipeline_options.do_table_structure = False

            converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(
                        pipeline_options=pipeline_options
                    )
                }
            )
            result   = converter.convert(str(src))
            doc      = result.document
            doc_dict = doc.export_to_dict()
        except Exception as e:
            logger.warning("Docling 解析失敗（%s）：%s", suffix, e)
            return None

        content_list: list[dict] = []
        cnt = 0

        def _walk(block: dict, btype: str) -> None:
            nonlocal cnt
            children = block.get("children")
            if not children:
                cnt += 1
                content_list.append(_docling_block_to_item(block, btype, img_dir, cnt))
            else:
                if btype not in ("groups", "body"):
                    cnt += 1
                    content_list.append(
                        _docling_block_to_item(block, btype, img_dir, cnt)
                    )
                for ref in children:
                    tag   = ref.get("$ref", "")
                    parts = tag.split("/")
                    if len(parts) < 3:
                        continue
                    mtype, mnum = parts[1], parts[2]
                    try:
                        _walk(doc_dict[mtype][int(mnum)], mtype)
                    except (KeyError, ValueError, IndexError):
                        pass

        if "body" in doc_dict:
            _walk(doc_dict["body"], "body")

        return content_list if content_list else None


def _docling_block_to_item(block: dict, btype: str, img_dir: Path, cnt: int) -> dict:
    page_idx = cnt // 10
    if btype == "texts":
        label = block.get("label", "")
        if label == "formula":
            return {"type": "equation", "text": block.get("orig", ""),
                    "text_format": "unknown", "page_idx": page_idx}
        return {"type": "text", "text": block.get("orig", ""),
                "text_level": 1 if label in ("section_header", "title") else 0,
                "page_idx": page_idx}
    elif btype == "pictures":
        try:
            uri   = block.get("image", {}).get("uri", "")
            parts = uri.split(",", 1)
            b64   = parts[1] if len(parts) == 2 else parts[0]
            img_path = img_dir / f"img_{cnt}.png"
            img_path.write_bytes(base64.b64decode(b64))
            return {"type": "image", "img_path": str(img_path.resolve()),
                    "image_caption": block.get("caption", ""), "page_idx": page_idx}
        except Exception:
            return {"type": "text",
                    "text": f"[圖片: {block.get('caption', '')}]",
                    "page_idx": page_idx}
    else:
        return {"type": "table", "table_body": block.get("data", []),
                "table_caption": block.get("caption", ""), "page_idx": page_idx}


# ══════════════════════════════════════════════════════════
# PaddleOCR 引擎（PaddleX 3.x 需要連網下載模型，暫時停用）
# 若需啟用，請確認模型已下載並設定 PDF_PARSER=paddleocr
# ══════════════════════════════════════════════════════════

def _check_paddleocr() -> bool:
    # PaddleX 3.x 架構需要從官方平台下載模型，網路受限環境無法使用
    # 若模型已在本地，可將此處改為 True 並確認 _run_paddleocr 正常運作
    return False



def _run_paddleocr(pdf_bytes: bytes) -> list[dict] | None:
    """PaddleX 3.x 需要連網下載模型，目前停用，直接回傳 None"""
    logger.debug("PaddleOCR 已停用（PaddleX 3.x 需連網），跳過")
    return None


# ══════════════════════════════════════════════════════════
# pdfplumber 引擎（最輕量 fallback）
# ══════════════════════════════════════════════════════════

def _run_pdfplumber(pdf_bytes: bytes) -> list[ParsedPage]:
    """兩遍掃描：第一遍收文字供 context window，第二遍完整萃取"""
    hash_counts: dict[str, int] = {}
    all_texts  : list[str]      = []

    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for page in pdf.pages:
            all_texts.append((page.extract_text() or "").replace("\x00", "").strip())
            for img in (page.images or []):
                try:
                    crop = page.crop((img["x0"], img["top"], img["x1"], img["bottom"]))
                    buf  = io.BytesIO()
                    crop.to_image().save(buf, format="PNG")
                    h = _md5(buf.getvalue())
                    hash_counts[h] = hash_counts.get(h, 0) + 1
                except Exception:
                    pass
    dup_hashes = {h for h, n in hash_counts.items() if n > 1}

    seen_img: dict[str, str] = {}
    pages   : list[ParsedPage] = []

    with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
        for idx, page in enumerate(pdf.pages):
            page_num = idx + 1
            raw = all_texts[idx]

            raw = re.sub(r"[^\n]*第\s*\d+\s*頁[^\n]*\n?", "", raw)
            raw = re.sub(r"本文件之內容為.+?其一部或全部。\n?", "", raw, flags=re.DOTALL)
            raw = re.sub(
                r"The information contained in this material.+?permission of \w+\.\n?",
                "", raw, flags=re.DOTALL | re.IGNORECASE,
            )
            raw = raw.strip()

            tables    = page.extract_tables() or []
            tbl_mds   : list[str] = []
            tbl_texts : set[str]  = set()
            for tbl in tables:
                md = _table_to_md(tbl)
                if md:
                    tbl_mds.append(md)
                for row in tbl:
                    for cell in (row or []):
                        if cell:
                            s = str(cell).strip()
                            if len(s) >= settings.TBL_DEDUP_MIN_LEN:
                                tbl_texts.add(s)
                                tbl_texts.add(re.sub(r"\s+", " ", s))
            if tbl_texts:
                raw = "\n".join(
                    ln for ln in raw.splitlines()
                    if re.sub(r"\s+", " ", ln.strip()) not in tbl_texts
                ).strip()

            ctx = _build_context_from_texts(all_texts, idx)

            fig_descs: list[str] = []
            for i, img in enumerate(page.images or []):
                try:
                    crop = page.crop((img["x0"], img["top"], img["x1"], img["bottom"]))
                    buf  = io.BytesIO()
                    crop.to_image().save(buf, format="PNG")
                    img_data = buf.getvalue()
                    if len(img_data) < settings.IMG_MIN_BYTES:
                        continue
                    h = _md5(img_data)
                    if h in dup_hashes:
                        continue
                    if h in seen_img:
                        fig_descs.append(seen_img[h])
                        continue
                    desc = _describe_with_context(img_data, page_num, ctx)
                    seen_img[h] = desc
                    fig_descs.append(desc)
                except Exception as e:
                    logger.warning("第 %d 頁圖片失敗：%s", page_num, e)

            if _is_scanned_page(raw) and not tbl_mds:
                try:
                    pg_img = page.to_image(resolution=150)
                    buf    = io.BytesIO()
                    pg_img.save(buf, format="PNG")
                    raw = _describe_with_context(buf.getvalue(), page_num, ctx)
                except Exception as e:
                    logger.warning("第 %d 頁整頁 Vision 失敗：%s", page_num, e)

            if fig_descs and not raw and not tbl_mds:
                ctype = ContentType.IMAGE
            elif tbl_mds and not raw and not fig_descs:
                ctype = ContentType.TABLE
            else:
                ctype = ContentType.MIXED

            pages.append(ParsedPage(
                page_num=page_num, text=raw,
                tables=tbl_mds, figure_desc=fig_descs, content_type=ctype,
            ))
    return pages


# ══════════════════════════════════════════════════════════
# 統一 PDF 解析入口
# ══════════════════════════════════════════════════════════

def _parse_as_pdf(pdf_bytes: bytes) -> list[ParsedPage]:
    """auto 模式：docling → pdfplumber"""
    parser = settings.PDF_PARSER.lower()

    def _try_docling():
        cl = _run_docling(pdf_bytes, ".pdf")
        return _content_list_to_pages(cl) if cl is not None else None

    def _try_paddleocr():
        cl = _run_paddleocr(pdf_bytes)
        return _content_list_to_pages(cl) if cl is not None else None

    def _try_plumber():
        return _run_pdfplumber(pdf_bytes)

    engine_map = {
        "docling"   : [("docling",    _try_docling)],
        "paddleocr" : [("pdfplumber", _try_plumber)],   # PaddleX 3.x 需連網，改用 pdfplumber
        "pdfplumber": [("pdfplumber", _try_plumber)],
        "auto"      : [("docling",    _try_docling),
                       ("pdfplumber", _try_plumber)],
    }
    engines = engine_map.get(parser, engine_map["auto"])

    for name, fn in engines:
        try:
            result = fn()
            if result is not None:
                logger.info("PDF 解析引擎：%s", name)
                return result
        except Exception as e:
            logger.warning("%s 失敗，嘗試下一個：%s", name, e)

    return _run_pdfplumber(pdf_bytes)


# ══════════════════════════════════════════════════════════
# 各格式公開解析函式
# ══════════════════════════════════════════════════════════

def parse_pdf(file_bytes: bytes) -> list[ParsedPage]:
    logger.info("解析 PDF（%d bytes）引擎=%s", len(file_bytes), settings.PDF_PARSER)
    return _parse_as_pdf(file_bytes)


def parse_office(file_bytes: bytes, suffix: str) -> list[ParsedPage]:
    logger.info("解析 Office（%s，%d bytes）", suffix, len(file_bytes))

    # 1. Docling 直接解析
    if _check_docling():
        cl = _run_docling(file_bytes, suffix)
        if cl is not None:
            logger.info("Docling 直接解析 %s", suffix)
            return _content_list_to_pages(cl)
        logger.warning("Docling 解析 %s 失敗，嘗試轉 PDF", suffix)

    # 2. 轉 PDF 再解析
    pdf_bytes, ok = _office_to_pdf(file_bytes, suffix)
    if ok:
        return _parse_as_pdf(pdf_bytes)

    # 3. 格式專用 fallback
    logger.warning("%s 無法轉 PDF，使用格式專用 fallback", suffix)
    return _office_native_fallback(file_bytes, suffix)


def _office_native_fallback(file_bytes: bytes, suffix: str) -> list[ParsedPage]:
    """格式專用 Python fallback（最後手段）"""
    s = suffix.lower()

    if s == ".doc":
        try:
            import docx2txt  # type: ignore
            with tempfile.NamedTemporaryFile(suffix=".doc", delete=False) as tmp:
                tmp.write(file_bytes); tmp_path = tmp.name
            try:
                text = docx2txt.process(tmp_path)
                if text and len(text.strip()) > 20:
                    return [ParsedPage(page_num=1, text=text.strip())]
            finally:
                try: os.unlink(tmp_path)
                except: pass
        except Exception:
            pass
        for enc in ("utf-16-le", "latin-1"):
            try:
                text = file_bytes.decode(enc, errors="ignore")
                text = re.sub(r"[^\x20-\x7E\u4e00-\u9fff\n\r\t]", " ", text)
                text = re.sub(r" {4,}", " ", text).strip()
                if len(text) > 50:
                    return [ParsedPage(page_num=1, text=text)]
            except Exception:
                pass
        return [ParsedPage(page_num=1, text="")]

    if s == ".docx":
        try:
            from docx import Document
            doc       = Document(io.BytesIO(file_bytes))
            lines     : list[str] = []
            tbl_mds   : list[str] = []
            fig_descs : list[str] = []
            seen_img  : dict[str, str] = {}
            doc_rels  : dict[str, bytes] = {}
            for rid, rel in doc.part.rels.items():
                if "image" not in rel.reltype:
                    continue
                try:
                    img_data = rel.target_part.blob
                    ct = rel.target_part.content_type
                    if "emf" in ct or "wmf" in ct or len(img_data) < settings.IMG_MIN_BYTES:
                        continue
                    doc_rels[rid] = img_data
                except Exception:
                    pass

            def _para_rids(para) -> list[str]:
                rids = []
                try:
                    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
                    for blip in para._element.findall(f".//{{{ns}}}blip"):
                        rid = blip.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                        if rid and rid in doc_rels:
                            rids.append(rid)
                    xml = para._element.xml
                    if isinstance(xml, bytes):
                        xml = xml.decode("utf-8", errors="replace")
                    for rid in re.findall(r'<v:imagedata\s+r:id="([^"]+)"', xml):
                        if rid in doc_rels:
                            rids.append(rid)
                except Exception:
                    pass
                return rids

            for para in doc.paragraphs:
                for rid in _para_rids(para):
                    img_data = doc_rels[rid]
                    h = _md5(img_data)
                    if h not in seen_img:
                        seen_img[h] = describe_image(img_data, context="Word 內嵌圖片")
                    lines.append(f"> **圖片描述**\n{seen_img[h]}")
                    fig_descs.append(seen_img[h])
                t = para.text.strip()
                if not t:
                    continue
                name  = (para.style.name or "") if para.style else ""
                level = 0
                if "Heading" in name:
                    try: level = int("".join(c for c in name if c.isdigit()) or "1")
                    except: level = 1
                lines.append("#" * level + " " + t if level > 0 else t)
            for tbl in doc.tables:
                rows = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
                md   = _table_to_md(rows)
                if md:
                    tbl_mds.append(md)
            full = re.sub(r"\n{3,}", "\n\n", "\n".join(lines))
            return [ParsedPage(page_num=1, text=full, tables=tbl_mds, figure_desc=fig_descs)]
        except Exception as e:
            logger.warning("python-docx fallback 失敗：%s", e)
            return [ParsedPage(page_num=1, text="")]

    if s == ".xls":
        try:
            import xlrd  # type: ignore
            wb    = xlrd.open_workbook(file_contents=file_bytes)
            pages : list[ParsedPage] = []
            for si in range(wb.nsheets):
                ws   = wb.sheet_by_index(si)
                rows = []
                for r in range(ws.nrows):
                    row = [str(ws.cell_value(r, c)).strip() for c in range(ws.ncols)]
                    if any(row):
                        rows.append(row)
                tbl_md = _table_to_md(rows) if rows else ""
                plain  = "\n".join("\t".join(r) for r in rows)
                pages.append(ParsedPage(si + 1, plain,
                                        tables=[tbl_md] if tbl_md else [],
                                        content_type=ContentType.TABLE))
            return pages if pages else [ParsedPage(page_num=1, text="")]
        except Exception as e:
            logger.warning("xlrd fallback 失敗：%s", e)
            return [ParsedPage(page_num=1, text="")]

    if s == ".xlsx":
        try:
            import openpyxl
            wb       = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            pages    : list[ParsedPage] = []
            seen_img : dict[str, str]   = {}
            for si, ws in enumerate(wb.worksheets):
                rows = []
                for row in ws.iter_rows(values_only=True):
                    c = [str(x).strip() if x is not None else "" for x in row]
                    if any(c):
                        rows.append(c)
                tbl_md    = _table_to_md(rows) if rows else ""
                plain     = "\n".join("\t".join(r) for r in rows)
                fig_descs : list[str] = []
                if hasattr(ws, "_images") and ws._images:
                    for img_obj in ws._images:
                        img_data = _safe_read_image_ref(img_obj)
                        if img_data:
                            h = _md5(img_data)
                            if h not in seen_img:
                                seen_img[h] = describe_image(img_data, context=f"Excel {ws.title}")
                            fig_descs.append(seen_img[h])
                pages.append(ParsedPage(si + 1, plain,
                                        tables=[tbl_md] if tbl_md else [],
                                        figure_desc=fig_descs,
                                        content_type=ContentType.TABLE))
            wb.close()
            return pages if pages else [ParsedPage(page_num=1, text="")]
        except Exception as e:
            logger.warning("openpyxl fallback 失敗：%s", e)
            return [ParsedPage(page_num=1, text="")]

    if s == ".ppt":
        for enc in ("utf-16-le", "latin-1"):
            try:
                text = file_bytes.decode(enc, errors="ignore")
                text = re.sub(r"[^\x20-\x7E\u4e00-\u9fff\n\r\t]", " ", text)
                text = re.sub(r" {4,}", " ", text).strip()
                if len(text) > 50:
                    return [ParsedPage(page_num=1, text=text)]
            except Exception:
                pass
        return [ParsedPage(page_num=1, text="")]

    if s == ".pptx":
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            prs      = Presentation(io.BytesIO(file_bytes))
            pages    : list[ParsedPage] = []
            seen_img : dict[str, str]   = {}
            all_slide_text: list[str]   = []
            for slide in prs.slides:
                parts = [para.text.strip()
                         for shape in slide.shapes
                         if shape.has_text_frame
                         for para in shape.text_frame.paragraphs
                         if para.text.strip()]
                all_slide_text.append("\n".join(parts))

            def _proc(shape, slide_num, parts_text, tbl_mds, fig_descs):
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for child in shape.shapes:
                        _proc(child, slide_num, parts_text, tbl_mds, fig_descs)
                    return
                if shape.has_table:
                    md = _table_to_md([[_clean_cell(c.text) for c in r.cells]
                                       for r in shape.table.rows])
                    if md:
                        tbl_mds.append(md)
                    return
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or hasattr(shape, "image"):
                    try:
                        img_data = shape.image.blob
                        ct       = shape.image.content_type
                        if "emf" not in ct and "wmf" not in ct \
                                and len(img_data) >= settings.IMG_MIN_BYTES:
                            h = _md5(img_data)
                            if h not in seen_img:
                                ctx = _build_context_from_texts(all_slide_text, slide_num - 1)
                                seen_img[h] = _describe_with_context(img_data, slide_num, ctx)
                            fig_descs.append(seen_img[h])
                    except Exception:
                        pass
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            lv = para.level or 0
                            parts_text.append("  " * (lv - 1) + "- " + t if lv > 0 else t)

            for si, slide in enumerate(prs.slides):
                slide_num  = si + 1
                parts_text : list[str] = []
                tbl_mds    : list[str] = []
                fig_descs  : list[str] = []
                for shape in sorted(slide.shapes, key=lambda s: (s.top or 0, s.left or 0)):
                    _proc(shape, slide_num, parts_text, tbl_mds, fig_descs)
                full = re.sub(r"\n{3,}", "\n\n", "\n".join(parts_text)).strip()
                pages.append(ParsedPage(slide_num, full, tbl_mds, fig_descs))
            return pages
        except Exception as e:
            logger.warning("python-pptx fallback 失敗：%s", e)
            return [ParsedPage(page_num=1, text="")]

    return [ParsedPage(page_num=1, text="")]


def parse_txt(file_bytes: bytes, suffix: str = ".txt") -> list[ParsedPage]:
    logger.info("解析 TXT/MD（%d bytes）", len(file_bytes))
    parser = settings.PDF_PARSER.lower()
    if parser != "pdfplumber":
        pdf_bytes = _txt_to_pdf(file_bytes)
        if pdf_bytes:
            return _parse_as_pdf(pdf_bytes)
    text = file_bytes.decode("utf-8", errors="replace")
    return [ParsedPage(page_num=1, text=text, content_type=ContentType.TEXT)]


def parse_image(file_bytes: bytes, filename: str = "") -> list[ParsedPage]:
    logger.info("解析圖片：%s", filename)
    desc = describe_image(
        file_bytes,
        context=f"圖片：{filename}，請完整描述所有文字與內容",
    )
    return [ParsedPage(page_num=1, text="", figure_desc=[desc],
                       content_type=ContentType.IMAGE)]


# ══════════════════════════════════════════════════════════
# 統一入口
# ══════════════════════════════════════════════════════════

_OFFICE_SUFFIXES = {".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx"}
_TEXT_SUFFIXES   = {".txt", ".md"}


def parse_file(file_bytes: bytes, filename: str) -> list[ParsedPage]:
    """唯一對外的公開函式"""
    from app.utils.helpers import get_file_type
    suffix = Path(filename).suffix.lower()
    ftype  = get_file_type(filename)

    if ftype == "image":
        return parse_image(file_bytes, filename)
    if ftype == "pdf":
        return parse_pdf(file_bytes)
    if suffix in _OFFICE_SUFFIXES:
        return parse_office(file_bytes, suffix)
    if suffix in _TEXT_SUFFIXES:
        return parse_txt(file_bytes, suffix)

    raise ValueError(f"不支援的檔案格式：{suffix}")